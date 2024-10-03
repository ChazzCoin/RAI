import base64
import io
import os
from pptx import Presentation
from docx.opc.constants import RELATIONSHIP_TYPE as RT
import pdfplumber
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from files.save import DataSaver
from F import LIST
from dataset.intake.PDF_v1 import FPDF
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
from dataset.TextCleaner import text_to_chroma


def encode_image(image_data:bytes):
    return base64.b64encode(image_data).decode('utf-8')

def vision_extract_text(file_path):
    """
    Extracts text and images from PDF, PPTX, and DOCX files, performing OCR where necessary.

    Args:
        file_path (str): The path to the file.

    Returns:
        list of dict: A list of dictionaries containing 'text' and 'image' for each page or slide.
    """
    file_ext = os.path.splitext(file_path)[1].lower()

    page_texts = []

    if file_ext == '.pdf':
        try:
            # Convert PDF pages to images
            pages = convert_from_path(file_path)
        except Exception as e:
            print(f"Error converting PDF to images: {e}")
            return []

        for idx, page in enumerate(pages):
            try:
                if page.mode != 'RGB':
                    page = page.convert('RGB')
                # Perform OCR on the image
                text = pytesseract.image_to_string(page)
                # Prepare data for ChromaDB
                page_data = {
                    'text': text,
                    'image': page  # PIL Image object
                }
                page_texts.append(page_data)
                print(f"Processed Page {idx + 1}")
            except Exception as e:
                print(f"Error performing OCR on page {idx + 1}: {e}")
                page_texts.append({'text': '', 'image': None})
    elif file_ext == '.docx':
        try:
            doc = Document(file_path)
        except Exception as e:
            print(f"Error reading DOCX file: {e}")
            return []

        # Extract text from paragraphs
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        text = '\n'.join(full_text)

        # Extract images
        images = []
        for rel in doc.part.rels.values():
            if rel.reltype == RT.IMAGE:
                image_data = rel.target_part.blob
                images.append(image_data)
                # Perform OCR on the image
                img = Image.open(io.BytesIO(image_data))
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                img_text = pytesseract.image_to_string(img)
                text += '\n' + img_text

        page_data = {
            'text': text,
            'image': images  # List of image byte data
        }
        page_texts.append(page_data)
        print(f"Processed DOCX file")
    elif file_ext == '.pptx':
        try:
            prs = Presentation(file_path)
        except Exception as e:
            print(f"Error reading PPTX file: {e}")
            return []

        for idx, slide in enumerate(prs.slides):
            slide_text = ''
            # Extract text from slide shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text += shape.text + '\n'
            # Extract images from slide shapes
            images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_part = shape.image
                    image_data = image_part.blob
                    image_ext = image_part.ext.lower()
                    images.append(image_data)
                    # Perform OCR on the image if it's a supported format
                    try:
                        if image_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff']:
                            img = Image.open(io.BytesIO(image_data))
                            # Handle palette images with transparency
                            if img.mode == 'P':
                                img = img.convert('RGBA')
                            elif img.mode != 'RGB' and img.mode != 'RGBA':
                                img = img.convert('RGB')
                            img_text = pytesseract.image_to_string(img)
                            slide_text += '\n' + img_text
                        else:
                            print(f"Unsupported image format {image_ext} in slide {idx + 1}")
                    except Exception as e:
                        print(f"Error performing OCR on an image in slide {idx + 1}: {e}")
            page_data = {
                'text': slide_text,
                'image': images  # List of image byte data
            }
            page_texts.append(page_data)
            print(f"Processed Slide {idx + 1}")

    else:
        print(f"Unsupported file type: {file_ext}")
        return []

    return page_texts

def vision_extract_text_from_pdf(pdf_path):
    """
    Extracts text from a PDF file by converting each page to an image and performing OCR.

    Args:
        pdf_path (str): The path to the PDF file.

    Returns:
        list of str: A list containing the text extracted from each page.
    """
    try:
        # Convert PDF pages to images
        pages = convert_from_path(pdf_path)
    except Exception as e:
        print(f"Error converting PDF to images: {e}")
        return []

    page_texts = []
    chroma_docs = []

    for idx, page in enumerate(pages):
        try:
            if page.mode != 'RGB':
                page = page.convert('RGB')
            # Perform OCR on the image
            text = pytesseract.image_to_string(page)
            page_texts.append(text)

            print(f"Raw Page {idx}: {text}")
            prepped = text_to_chroma(text)
            chroma_docs.append(prepped)
            chroma_docs = LIST.flatten(chroma_docs)
            print(f"Chroma Page {idx}: {prepped}")
        except Exception as e:
            print(f"Error performing OCR on page {idx + 1}: {e}")
            page_texts.append("")

    return page_texts

# def to_chroma_docs(extracted_data):
#     # Process extracted_data to save into ChromaDB
#     docs = []
#     for item in extracted_data:
#         text = item['text']
#         image = item['image']
#
#         # Generate text embeddings
#         text_embedding = text_embedding_model.embed(text)
#
#         # Generate image embeddings
#         if image:
#             if isinstance(image, list):
#                 # For DOCX and PPTX (list of images)
#                 image_embeddings = [image_embedding_model.embed(Image.open(io.BytesIO(img_data))) for img_data in image]
#             else:
#                 # For PDFs (single PIL Image)
#                 image_embedding = image_embedding_model.embed(image)
#
#         # Prepare data for ChromaDB
#         chroma_document = {
#             'text_embedding': text_embedding,
#             'image_embedding': image_embedding if not isinstance(image, list) else image_embeddings,
#             'metadata': {
#                 'source_file': file_path
#             }
#         }
#         docs.append(chroma_document)
#         # Insert into ChromaDB
#         # chromadb_client.insert(chroma_document)
#     return docs

"""
    ------------
"""

def pdf_to_txt_file(pdf_file, output_text_file:str):
    """
    Extracts text from a PDF document and saves it to a text file.
    Args:
        pdf_file (str): The path to the PDF document.
        output_text_file (str): The path to the output text file where extracted data will be saved.
    """
    try:
        # Initialize an empty string to collect the extracted text
        extracted_text = []
        # Open the PDF file
        with pdfplumber.open(pdf_file) as pdf:
            # Loop through each page of the PDF
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    extracted_text.append(f"--- Page {page_num + 1} ---\n{text}\n")
        DataSaver.save_txt(data=extracted_text, output=output_text_file)
        print(f"Text successfully extracted to {output_text_file}")
    except Exception as e:
        print(f"Error processing the PDF document {pdf_file}: {e}")

def pdf_to_txt(pdf_file):
    """
    Extracts text from a PDF document and saves it to a text file.
    Args:
        pdf_file (str): The path to the PDF document.
        output_text_file (str): The path to the output text file where extracted data will be saved.
    """
    try:
        # Initialize an empty string to collect the extracted text
        extracted_text = []
        # Open the PDF file
        with pdfplumber.open(pdf_file) as pdf:
            # Loop through each page of the PDF
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    extracted_text.append(f"--- Page {page_num + 1} ---\n{text}\n")
        print(f"Text successfully extracted pdf.")
        return LIST.to_str(extracted_text)
    except Exception as e:
        print(f"Error processing the PDF document {pdf_file}: {e}")
        return FPDF.extract_text_from_pdf(pdf_file)



if __name__ == "__main__":
    # from dataset.intake import PowerPoint
    # Example usage
    # pdf_file_path = "/Users/chazzromeo/Desktop/ussfTrainingData/us-soccer-player-development-framework-age-group-learning-plans.pdf"
    pdf_file_path = "/Users/chazzromeo/Desktop/ParkCityTrainingData/PCSCIndividualDevelopmentPlan11U12U.docx"
    # pdf_file = PowerPoint.convert_pptx_to_pdf(pdf_file_path)
    # output_text_file_path = f"{default_file_path()}/PCSC Player Handbook.txt"

    # pdf_to_txt_file(pdf_file_path, output_text_file_path)
    # text = pdf_to_txt(pdf_file_path)
    vision = vision_extract_text(pdf_file_path)
    print(f"vision: {len(vision)}")

    image = Image.open(io.BytesIO(vision[0]['image'][0]))

    # Display the image (this opens the image using the default image viewer)
    image.show()

    # Optionally, save the image to a file
    image.save('/Users/chazzromeo/Desktop/output_image.png')
    # Image.open(vision[0]['image'][0])