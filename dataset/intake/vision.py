import base64
import io
import os
import shutil
import stat

from pptx import Presentation
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
from F.LOG import Log
from F import OS


Log = Log("VisionExtractor")

def encode_image(image_data:bytes):
    return base64.b64encode(image_data).decode('utf-8')




class VisionExtractor:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def extract(self, save_images=False):
        return self._vision_extract_text(self.file_path, save_images)

    @staticmethod
    def _vision_extract_text(file_path, save_images):
        """
        Extracts text and images from PDF, PPTX, and DOCX files, performing OCR where necessary.

        Args:
            file_path (str): The path to the file.
            save_images (bool): Whether to save images to disk.

        Returns:
            list of dict: A list of dictionaries containing 'text' and 'image_path' or 'image' for each page or slide.
        """
        file_ext = os.path.splitext(file_path)[1].lower()

        page_texts = []
        output_dir = os.path.splitext(file_path)[0] + "_output"
        if save_images and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if file_ext == '.pdf':
            try:
                # Convert PDF pages to images
                pages = convert_from_path(file_path)
            except Exception as e:
                Log.e(f"Error converting PDF to images: {e}")
                return []
            for idx, page in enumerate(pages):
                try:
                    if page.mode != 'RGB':
                        page = page.convert('RGB')
                    # Perform OCR on the image
                    text = pytesseract.image_to_string(page)
                    # Save image if required
                    image_path = 'None'
                    if save_images:
                        image_path = os.path.join(output_dir, f"page_{idx + 1}.png")
                        page.save(image_path)
                    # Prepare data
                    page_data = {
                        'page_content': text,
                        'image_path': image_path if save_images else ''
                    }
                    page_texts.append(page_data)
                except Exception as e:
                    Log.e(f"Error performing OCR on page {idx + 1}: {e}")
                    page_texts.append({'page_content': '', 'image_path': ''})
        elif file_ext == '.docx':
            try:
                doc = Document(file_path)
            except Exception as e:
                print(f"Error reading DOCX file: {e}")
                return []

            # Extract text from paragraphs
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            text = '\n'.join(full_text)

            # Extract images
            images = []
            for rel_idx, rel in enumerate(doc.part.rels.values()):
                if rel.reltype == RT.IMAGE:
                    image_data = rel.target_part.blob
                    images.append(image_data)
                    # Save image if required
                    if save_images:
                        image_path = os.path.join(output_dir, f"image_{rel_idx + 1}.png")
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_data)
                    # Perform OCR on the image
                    img = Image.open(io.BytesIO(image_data))
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    img_text = pytesseract.image_to_string(img)
                    text += '\n' + img_text
            page_data = {
                'page_content': text,
                'image_paths': [os.path.join(output_dir, f"image_{i + 1}.png") for i in range(len(images))] if save_images else images
            }
            page_texts.append(page_data)
        elif file_ext == '.pptx':
            try:
                prs = Presentation(file_path)
            except Exception as e:
                Log.e(f"Error reading PPTX file: {e}")
                return []

            for idx, slide in enumerate(prs.slides):
                slide_text = ''
                # Extract text from slide shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text += shape.text + '\n'
                # Extract images from slide shapes
                images = []
                image_paths = []
                for img_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_part = shape.image
                        image_data = image_part.blob
                        image_ext = image_part.ext.lower()
                        images.append(image_data)
                        # Save image if required
                        if save_images:
                            image_path = os.path.join(output_dir, f"slide_{idx + 1}_image_{img_idx + 1}.png")
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image_data)
                            image_paths.append(image_path)
                        # Perform OCR on the image if it's a supported format
                        try:
                            if image_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff']:
                                img = Image.open(io.BytesIO(image_data))
                                # Handle palette images with transparency
                                if img.mode == 'P':
                                    img = img.convert('RGBA')
                                elif img.mode != 'RGB' and img.mode != 'RGBA':
                                    img = img.convert('RGB')
                                img_text = pytesseract.image_to_string(img)
                                slide_text += '\n' + img_text
                            else:
                                Log.w(f"Unsupported image format {image_ext} in slide {idx + 1}")
                        except Exception as e:
                            Log.e(f"Error performing OCR on an image in slide {idx + 1}: {e}")
                page_data = {
                    'page_content': slide_text,
                    'image_paths': image_paths if save_images else images  # List of image paths or byte data
                }
                page_texts.append(page_data)
        else:
            Log.e(f"Unsupported file type: {file_ext}")
            return []

        Log.s(f"Successfully extracted text and images from: [ {file_path} ]")
        return page_texts

