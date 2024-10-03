from pptx import Presentation
from config import default_file_path
from files.FilePath import FilePath
from files.save import DataSaver
import subprocess
import os
import platform
import shutil

def convert_pptx_to_pdf(pptx_path, output_dir=None):
    """
    Converts a .pptx PowerPoint file to a PDF file using LibreOffice in headless mode.

    Args:
        pptx_path (str): The path to the .pptx file to be converted.
        output_dir (str, optional): The directory where the PDF file will be saved.
                                    If not provided, uses the same directory as the pptx file.

    Returns:
        str: The path to the converted PDF file.

    Raises:
        FileNotFoundError: If the input .pptx file does not exist.
        RuntimeError: If the conversion fails.
    """
    if not os.path.isfile(pptx_path):
        raise FileNotFoundError(f"The file {pptx_path} does not exist.")

    # Determine the output directory
    if output_dir is None:
        output_dir = os.path.dirname(pptx_path)
    else:
        os.makedirs(output_dir, exist_ok=True)

    # Construct the command based on the operating system
    system_platform = platform.system()
    if system_platform == 'Windows':
        # Find the LibreOffice installation path
        program_files = os.environ.get('ProgramFiles', r'C:\Program Files')
        libreoffice_paths = [
            os.path.join(program_files, 'LibreOffice', 'program', 'soffice.exe'),
            os.path.join(program_files + ' (x86)', 'LibreOffice', 'program', 'soffice.exe')
        ]
        libreoffice_exec = None
        for path in libreoffice_paths:
            if os.path.isfile(path):
                libreoffice_exec = path
                break
        if not libreoffice_exec:
            raise FileNotFoundError("LibreOffice executable not found. Please ensure LibreOffice is installed.")
    else:
        # For macOS and Linux, assume 'soffice' is in PATH
        libreoffice_exec = shutil.which('soffice')
        if not libreoffice_exec:
            raise FileNotFoundError("LibreOffice executable not found in PATH. Please ensure LibreOffice is installed.")

    # Prepare the command to convert pptx to pdf
    cmd = [
        libreoffice_exec,
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', output_dir,
        pptx_path
    ]

    try:
        # Execute the command
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Conversion failed: {e}")

    # Construct the output PDF file path
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_file = os.path.join(output_dir, base_name + '.pdf')

    if not os.path.isfile(pdf_file):
        raise RuntimeError("Conversion failed: PDF file was not created.")

    return pdf_file


def pptx_to_txt_file(pptx_file, output_text_file:FilePath):
    # Open the presentation
    presentation = Presentation(pptx_file)
    # Initialize an empty list to hold the extracted text
    text_runs = []
    # Loop through each slide in the presentation
    for slide in presentation.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape contains text
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    # Write the extracted text to the output text file
    DataSaver.save_txt(data=text_runs, output=output_text_file)
    print(f"Text successfully extracted to {output_text_file}")

if __name__ == "__main__":
    # Example usage
    pptx_file_path = "/Users/chazzromeo/Desktop/parkcitysoccer/PCSC Pre-Placement May 5, 2024.pptx"
    output_text_file_path = f"{default_file_path()}/pptx-test.txt"
    pptx_to_txt_file(pptx_file_path, output_text_file_path)
