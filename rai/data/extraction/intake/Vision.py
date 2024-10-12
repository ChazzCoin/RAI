import base64
import io
import os
from pptx import Presentation
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
from F.LOG import Log
from rai.data import RaiPath
Log = Log("VisionExtractor")

def encode_image(image_data:bytes):
    return base64.b64encode(image_data).decode('utf-8')

class VisionExtractor:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def extract(self, save_images=False):
        data = self._vision_extract_text(self.file_path, save_images)
        if not self.validate_vision_results(data):
            data = None
        return data

    @staticmethod
    def _vision_extract_text(file_path, save_images):
        """
        Extracts text and images from PDF, PPTX, and DOCX files, performing OCR where necessary.

        Args:
            file_path (str): The path to the file.
            save_images (bool): Whether to save images to disk.

        Returns:
            list of dict: A list of dictionaries containing 'text' and 'image_path' or 'image' for each page or slide.
        """
        p = RaiPath(file_path)
        file_ext = p.ext_type

        page_texts = []
        output_dir = os.path.splitext(file_path)[0] + "_output"
        if save_images and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        try:
            if file_ext in ['pdf', '.pdf']:
                try:
                    # Convert PDF pages to images
                    pages = convert_from_path(file_path)
                except Exception as e:
                    Log.w("Error opening pdf.", e)
                    return None
                for idx, page in enumerate(pages):
                    try:
                        if page.mode != 'RGB':
                            page = page.convert('RGB')
                        # Perform OCR on the image
                        text = pytesseract.image_to_string(page)
                        # Save image if required
                        image_path = 'None'
                        if save_images:
                            image_path = os.path.join(output_dir, f"page_{idx + 1}.png")
                            page.save(image_path)
                        # Prepare data
                        page_data = {
                            'page_content': text,
                            'image_path': image_path if save_images else ''
                        }
                        page_texts.append(page_data)
                    except Exception as e:
                        Log.w(f"Error performing OCR on page {idx + 1}: {e}")
                        page_texts.append({'page_content': '', 'image_path': ''})
                        continue
            elif file_ext in ['docx', 'docx']:
                try:
                    doc = Document(file_path)
                except Exception as e:
                    Log.w("Error opening word doc.", e)
                    return None

                # Extract text from paragraphs
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                text = '\n'.join(full_text)

                # Extract images
                images = []
                for rel_idx, rel in enumerate(doc.part.rels.values()):
                    try:
                        if rel.reltype == RT.IMAGE:
                            image_data = rel.target_part.blob
                            images.append(image_data)
                            # Save image if required
                            if save_images:
                                image_path = os.path.join(output_dir, f"image_{rel_idx + 1}.png")
                                with open(image_path, 'wb') as img_file:
                                    img_file.write(image_data)
                            # Perform OCR on the image
                            img = Image.open(io.BytesIO(image_data))
                            if img.mode != 'RGB':
                                img = img.convert('RGB')
                            img_text = pytesseract.image_to_string(img)
                            text += '\n' + img_text
                        page_data = {
                            'page_content': text,
                            'image_paths': [os.path.join(output_dir, f"image_{i + 1}.png") for i in
                                            range(len(images))] if save_images else images
                        }
                        page_texts.append(page_data)
                    except Exception as e:
                        Log.w(f"Error performing OCR on page: {rel_idx}", e)
                        continue
            elif file_ext in ['pptx', '.pptx']:
                try:
                    prs = Presentation(file_path)
                except Exception as e:
                    Log.w("Error opening powerpoint.", e)
                    return None

                for idx, slide in enumerate(prs.slides):
                    try:
                        slide_text = ''
                        # Extract text from slide shapes
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                slide_text += shape.text + '\n'
                        # Extract images from slide shapes
                        images = []
                        image_paths = []
                        for img_idx, shape in enumerate(slide.shapes):
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                image_part = shape.image
                                image_data = image_part.blob
                                image_ext = image_part.ext.lower()
                                images.append(image_data)
                                # Save image if required
                                if save_images:
                                    image_path = os.path.join(output_dir, f"slide_{idx + 1}_image_{img_idx + 1}.png")
                                    with open(image_path, 'wb') as img_file:
                                        img_file.write(image_data)
                                    image_paths.append(image_path)
                                # Perform OCR on the image if it's a supported format
                                try:
                                    if image_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff']:
                                        img = Image.open(io.BytesIO(image_data))
                                        # Handle palette images with transparency
                                        if img.mode == 'P':
                                            img = img.convert('RGBA')
                                        elif img.mode != 'RGB' and img.mode != 'RGBA':
                                            img = img.convert('RGB')
                                        img_text = pytesseract.image_to_string(img)
                                        slide_text += '\n' + img_text
                                    else:
                                        Log.w(f"Unsupported image format {image_ext} in slide {idx + 1}")
                                except Exception as e:
                                    Log.w(f"Error performing OCR on an image in slide {idx + 1}: {e}")
                        page_data = {
                            'page_content': slide_text,
                            'image_paths': image_paths if save_images else images  # List of image paths or byte data
                        }
                        page_texts.append(page_data)
                    except Exception as e:
                        Log.w(f"Error performing OCR on a slide {idx + 1}: {e}")
                        continue
            else:
                Log.e(f"Unsupported file type: {file_ext}")
                return page_texts

            Log.s(f"Successfully extracted text and images from: [ {file_path} ]")
            # p.change_ext_type('txt')
            # save_list_to_file_with_headers(page_texts, f"/Users/chazzromeo/Desktop/visionout/{p.file_name}")
            return page_texts
        except Exception as e:
            Log.w(f"Error extracting text and images from: [ {file_path} ]", e)
            if page_texts is None:
                page_texts = []
            return page_texts

    @staticmethod
    def validate_vision_results(page_texts: list):
        """
        Validates that the page_texts list contains valid data.

        The list should:
        - Not be None
        - Not be empty
        - Contain no empty objects within it
        - Contain objects that match the structure {'page_content': '', 'image_path': ''}

        Args:
            page_texts (list): The list to validate.

        Returns:
            bool: True if valid, False otherwise.
        """
        if page_texts is None:
            return False
        if not isinstance(page_texts, list):
            return False
        if len(page_texts) == 0:
            return False
        for index, item in enumerate(page_texts):
            if not isinstance(item, dict):
                return False
            if len(item) == 0:
                return False
            expected_keys = {'page_content', 'image_path'}
            if set(item.keys()) != expected_keys:
                return False
            if not isinstance(item['page_content'], str):
                return False
            if not isinstance(item['image_path'], str):
                return False
        return True

def save_list_to_file_with_headers(text_list, filename):
    try:
        with open(filename, 'w') as file:
            for i, text in enumerate(text_list, start=1):
                file.write(f"Page {i}\n")
                file.write(f"{text}\n")
                file.write("-" * 20 + "\n")  # Page separator
        # print(f"List successfully saved to {filename}")
    except Exception as e:
        print(f"An error occurred: {e}")