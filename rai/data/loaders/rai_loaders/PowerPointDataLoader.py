import os

from F.LOG import Log
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from langchain_core.document_loaders import BaseLoader
from pptx import Presentation

from rai.data.loaders import verify_loader_data
from rai.data.loaders.rai_loaders.LastResortDataLoader import LastResortDataLoader
from rai.data.loaders.rai_loaders.RaiLoaderDocument import RaiLoaderDocument
from rai.data.loaders.rai_loaders.VisionDataLoader import VisionDataLoader

Log = Log("PowerPointDataLoader")


class PowerPointDataLoader(BaseLoader):

    def __init__(self, file_path: str, metadata=None):
        if metadata is None:
            metadata = {'image': ''}
        self.file_path = file_path
        self.metadata = metadata
        self._validate_file_path()

    def _validate_file_path(self):
        if not os.path.isfile(self.file_path):
            raise FileNotFoundError(f"The file {self.file_path} does not exist.")
        if not self.file_path.endswith('.pptx'):
            raise ValueError("Unsupported file type. Only PowerPoint (.pptx) files are allowed.")

    def load(self) -> [str]:
        """
        Loads and extracts the text from the PowerPoint file (.pptx).
        The output is a list of formatted strings representing each slide,
        useful for embeddings, querying, and AI processing.
        """
        try:
            presentation = Presentation(self.file_path)
            formatted_slides = self.format_presentation(presentation)
            if formatted_slides:
                docs = []
                for slide in formatted_slides:
                    docs.append(RaiLoaderDocument(page_content=slide, metadata=self.metadata))
                return docs
            else:
                raise ValueError("Presentation is empty or could not be processed.")
        except Exception as e:
            # Fallback to VisionDataLoader if the primary method fails
            Log.w("PowerPointLoader Failed, falling back.", e)
            loader = VisionDataLoader(self.file_path)
            if loader.should_fallback():
                # Fallback to UnstructuredPowerPointLoader if VisionDataLoader is empty
                Log.w("Vision Failed, falling back.", e)
                loader = UnstructuredPowerPointLoader(self.file_path)
                if not verify_loader_data(loader):
                    # Fallback to LastResortLoader if all else fails
                    loader = LastResortDataLoader(self.file_path)
            return loader.load()

    @staticmethod
    def format_presentation(presentation: Presentation) -> [str]:
        """
        Formats the PowerPoint presentation into a list of strings,
        each representing a slide in a consistent and structured format.
        """
        Log.i("Formatting Presentation...")
        formatted_slides = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_texts.append(shape.text)
            slide_content = " | ".join(slide_texts)
            formatted_slides.append(f"Slide {slide_num}: {slide_content}")

        return formatted_slides