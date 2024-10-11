import json
import os
import shutil
import stat
import time

import PyPDF2
import docx
import jsonlines
import pandas as pd
from F import DICT
from F.LOG import Log
from langchain_community.document_loaders import (
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    PyPDFLoader,
    TextLoader,
    UnstructuredEPubLoader,
    UnstructuredExcelLoader,
    UnstructuredMarkdownLoader,
    UnstructuredPowerPointLoader,
    UnstructuredRSTLoader,
    UnstructuredXMLLoader,
)
from langchain_core.document_loaders import BaseLoader
from pptx import Presentation

from rai.data.intake.Vision import VisionExtractor
from rai.data.read import read_file
from rai.data import RaiPath
Log = Log("RaiDataLoaders")

DEFAULT_CONTENT = "This data is unknown at this time."
DEFAULT_METADATA = {
    'image':''
}

class RaiMetadataLoader:
    file:str
    metadata:dict = DEFAULT_METADATA

    def __init__(self, meta_file:str=None, meta_dict:dict=None):
        if meta_file: self.file = meta_file
        if meta_dict: self.metadata = meta_dict

    def load_from_meta_file(self):
        pass
    def get_metadata(self):
        return self.metadata if not None else DEFAULT_METADATA
    @staticmethod
    def default_metadata():
        return DEFAULT_METADATA


class RaiDataLoader:
    file:str = None
    loader: BaseLoader = None
    """ -> private DATA LOADER HELPER <- """

    def __init__(self, file):
        self.file = file
        self.loader = self.get_loader(file)

    @staticmethod
    def get_loader(file_path: str) -> BaseLoader:
        rai_file = RaiPath(file_path)
        file_ext = rai_file.ext_type
        Log.w(f"Parsing File: [ {rai_file.file_name} ] Ext: [ {file_ext} ] to get matching DataLoader...")
        try:
            if file_ext == "pdf":
                loader = PdfLoader(file_path)
            elif file_ext in ["csv", "xls", "xlsx"]:
                loader = TableLoader(file_path)
            elif file_ext == "jsonl":
                loader = JSONLLoader(file_path)
            elif file_ext == "json":
                loader = JSONLoader(file_path)
            elif file_ext == "rst":
                loader = UnstructuredRSTLoader(file_path, mode="elements")
            elif file_ext == "xml":
                loader = UnstructuredXMLLoader(file_path)
            elif file_ext in ["htm", "html"]:
                loader = BSHTMLLoader(file_path, open_encoding="unicode_escape")
            elif file_ext == "md":
                loader = UnstructuredMarkdownLoader(file_path)
            elif file_ext == "epub":
                loader = UnstructuredEPubLoader(file_path)
            elif file_ext == "docx":
                loader = WordDocLoader(file_path)
            elif file_ext in ["ppt", "pptx"]:
                loader = PowerPointDataLoader(file_path)
            elif file_ext == "msg":
                loader = OutlookMessageLoader(file_path)
            elif file_ext == "txt":
                loader = LastResortLoader(file_path)
            elif file_ext in known_source_ext:
                loader = TextLoader(file_path, autodetect_encoding=True)
            else:
                loader = TextLoader(file_path, autodetect_encoding=True)
        except Exception as e:
            Log.w("Failed, falling back to original read file.", e)
            loader = None
        """ Fallback Intake Method """
        if not loader or not RaiDataLoader.verify_loader_data(loader):
            loader = LastResortLoader(file_path)
        return loader

    @staticmethod
    def verify_loader_data(loader: BaseLoader) -> bool:
        try:
            # Attempt to load data from the provided loader
            data = loader.load()
            # Check if data is not None, is a list, and has valid contents
            if data is None:
                return False
            if not isinstance(data, list):
                return False
            if len(data) == 0:
                return False
            # Check that each item in the data list is not None and has meaningful content
            for item in data:
                if item is None:
                    return False
                # Assuming the loader returns Document objects, check their content
                if not hasattr(item, 'page_content') or not isinstance(item.page_content, str):
                    return False
                if len(item.page_content.strip()) == 0:
                    return False
            # If all checks pass, return True
            return True
        except Exception as e:
            # If any exception occurs, log it if needed and return False
            # You could add logging here for better traceability in production
            Log.w("Failed Loader Validation", e)
            return False

    """ TODO: A METADATA LOADER 
        -> A METADATA FILE IN EACH FOLDER
        OR A DEFAULT METADATA OBJECT...
        EXTRACT METADATA FROM FILES
    """

"""
    CUSTOM LOADERS
"""
class RaiLoaderDocument:
    page_content:str
    metadata:dict

    def __init__(self, page_content: str, metadata:dict={ 'image':'' }) -> None:
        self.page_content = page_content
        self.metadata = metadata if not None else { 'image': '' }

class LastResortLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        """
        Attempts to read the file using a custom fallback method.
        If successful, it initializes a RawTextLoader with the extracted text.
        If everything fails, returns a valid data structure with a message indicating that the data is unknown.
        """
        try:
            raw_text = read_file(self.file_path)
            if raw_text:
                loader = RawTextLoader(raw_text)
                return loader.load()
            else:
                raise ValueError("Failed to extract raw text from the file.")
        except Exception as e:
            Log.t("FINAL DATA LOADER FAILURE ALERT! Falling Back: [ This data is unknown at this time. ]", e)
            return [RaiLoaderDocument(page_content=DEFAULT_CONTENT, metadata=DEFAULT_METADATA)]

class PdfLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path
        self._validate_file_path()

    def _validate_file_path(self):
        if not os.path.isfile(self.file_path):
            raise FileNotFoundError(f"The file {self.file_path} does not exist.")
        if not self.file_path.endswith('.pdf'):
            raise ValueError("Unsupported file type. Only PDF files are allowed.")

    def load(self) -> [str]:
        """
        Loads and extracts the text from the PDF file.
        The output is a list of formatted strings representing each page,
        useful for embeddings, querying, and AI processing.
        """
        try:
            with open(self.file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                formatted_pages = self.format_pdf(reader)
                if formatted_pages:
                    return formatted_pages
                else:
                    raise ValueError("PDF is empty or could not be processed.")
        except Exception as e:
            # Fallback to VisionDataLoader if the primary method fails
            Log.w("PDF Failed, falling back.", e)
            loader = VisionDataLoader(self.file_path)
            if loader.is_empty():
                # Fallback to PyPDFLoader if VisionDataLoader is empty
                loader = PyPDFLoader(self.file_path, extract_images=True)
                if not RaiDataLoader.verify_loader_data(loader):
                    # Fallback to LastResortLoader if all else fails
                    loader = LastResortLoader(self.file_path)
            return loader.load()

    @staticmethod
    def format_pdf(reader: PyPDF2.PdfReader) -> [str]:
        """
        Formats the PDF content into a list of strings,
        each representing a page in a consistent and structured format.
        """
        Log.i("Formatting Pdf...")
        formatted_pages = []
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            page_text = page.extract_text() if page.extract_text() else ""
            formatted_pages.append(f"Page {page_num + 1}: {page_text}")

        return formatted_pages

class TableLoader(BaseLoader):
    def __init__(self, file_path: str, metadata:dict={ 'image':'' }):
        self.file_path = file_path
        self.metadata = metadata if not None else { 'image': '' }
        self._validate_file_path()

    def _validate_file_path(self):
        if not os.path.isfile(self.file_path):
            raise FileNotFoundError(f"The file {self.file_path} does not exist.")
        if not (self.file_path.endswith('.csv') or self.file_path.endswith(('.xls', '.xlsx'))):
            raise ValueError("Unsupported file type. Only CSV and Excel files are allowed.")

    def load(self) -> [str]:
        """
        Loads and extracts the text from the table file (CSV or Excel).
        The output is a list of formatted strings representing each row,
        useful for embeddings, querying, and AI processing.
        """
        try:
            if self.file_path.endswith('.csv'):
                df = pd.read_csv(self.file_path)
            else:
                df = pd.read_excel(self.file_path)
            formatted_data = self.format_dataframe(df)
            if formatted_data:
                return formatted_data
            else:
                raise ValueError("Dataframe is empty or could not be processed.")
        except Exception as e:
            # Fallback to alternative loaders if the primary method fails
            Log.w("TableLoader Failed, falling back.", e)
            if self.file_path.endswith('.csv'):
                loader = CSVLoader(self.file_path)
            else:
                loader = UnstructuredExcelLoader(self.file_path)
            if not RaiDataLoader.verify_loader_data(loader):
                # Fallback to LastResortLoader if all else fails
                Log.w("TableLoader Failed, last resorting.", e)
                loader = LastResortLoader(self.file_path)
            return loader.load()

    @staticmethod
    def format_dataframe(df: pd.DataFrame) -> [str]:
        """
        Formats the DataFrame into a list of strings,
        each representing a row in a consistent and structured format.
        """
        Log.i("Formatting Dataframe...")
        formatted_rows = []
        headers = df.columns.tolist()
        header_str = " | ".join(headers)

        # Add headers as the first line to provide context for embeddings.
        formatted_rows.append(f"Headers: {header_str}")

        for index, row in df.iterrows():
            row_data = [f"{headers[i]}: {str(row[i])}" for i in range(len(headers))]
            row_str = " | ".join(row_data)
            formatted_rows.append(row_str)

        return formatted_rows

class WordDocLoader(BaseLoader):
    def __init__(self, file_path: str, metadata:dict=DEFAULT_METADATA):
        self.file_path = file_path
        self.metadata = metadata
        self._validate_file_path()

    def _validate_file_path(self):
        if not os.path.isfile(self.file_path):
            raise FileNotFoundError(f"The file {self.file_path} does not exist.")
        if not self.file_path.endswith('.docx'):
            raise ValueError("Unsupported file type. Only DOCX files are allowed.")

    def load(self) -> [str]:
        """
        Loads and extracts the text from the DOCX file.
        The output is a list of formatted strings representing each paragraph,
        useful for embeddings, querying, and AI processing.
        """
        try:
            document = docx.Document(self.file_path)
            formatted_paragraphs = self.format_docx(document)
            if formatted_paragraphs:
                return formatted_paragraphs
            else:
                raise ValueError("DOCX is empty or could not be processed.")
        except Exception as e:
            # Fallback to VisionDataLoader if the primary method fails
            loader = VisionDataLoader(self.file_path)
            if loader.is_empty():
                # Fallback to PyPDFLoader if VisionDataLoader is empty
                Log.w("Vision Failed, falling back.", e)
                loader = Docx2txtLoader(self.file_path)
                if not RaiDataLoader.verify_loader_data(loader):
                    # Fallback to LastResortLoader if all else fails
                    Log.w("WordDocLoader Failed, last resorting.", e)
                    loader = LastResortLoader(self.file_path)
            return loader.load()
    @staticmethod
    def format_docx(document: docx.Document) -> [str]:
        """
        Formats the DOCX content into a list of strings,
        each representing a paragraph in a consistent and structured format.
        """
        Log.i("Formatting Word Doc...")
        formatted_paragraphs = []
        for para_num, paragraph in enumerate(document.paragraphs, start=1):
            if paragraph.text.strip():
                formatted_paragraphs.append(f"Paragraph {para_num}: {paragraph.text.strip()}")

        return formatted_paragraphs

class PowerPointDataLoader(BaseLoader):

    def __init__(self, file_path: str, metadata=None):
        if metadata is None:
            metadata = {'image': ''}
        self.file_path = file_path
        self.metadata = metadata
        self._validate_file_path()

    def _validate_file_path(self):
        if not os.path.isfile(self.file_path):
            raise FileNotFoundError(f"The file {self.file_path} does not exist.")
        if not self.file_path.endswith('.pptx'):
            raise ValueError("Unsupported file type. Only PowerPoint (.pptx) files are allowed.")

    def load(self) -> [str]:
        """
        Loads and extracts the text from the PowerPoint file (.pptx).
        The output is a list of formatted strings representing each slide,
        useful for embeddings, querying, and AI processing.
        """
        try:
            presentation = Presentation(self.file_path)
            formatted_slides = self.format_presentation(presentation)
            if formatted_slides:
                return formatted_slides
            else:
                raise ValueError("Presentation is empty or could not be processed.")
        except Exception as e:
            # Fallback to VisionDataLoader if the primary method fails
            Log.w("PowerPointLoader Failed, falling back.", e)
            loader = VisionDataLoader(self.file_path)
            if loader.is_empty():
                # Fallback to UnstructuredPowerPointLoader if VisionDataLoader is empty
                Log.w("Vision Failed, falling back.", e)
                loader = UnstructuredPowerPointLoader(self.file_path)
                if not RaiDataLoader.verify_loader_data(loader):
                    # Fallback to LastResortLoader if all else fails
                    loader = LastResortLoader(self.file_path)
            return loader.load()

    @staticmethod
    def format_presentation(presentation: Presentation) -> [str]:
        """
        Formats the PowerPoint presentation into a list of strings,
        each representing a slide in a consistent and structured format.
        """
        Log.i("Formatting Presentation...")
        formatted_slides = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_texts.append(shape.text)
            slide_content = " | ".join(slide_texts)
            formatted_slides.append(f"Slide {slide_num}: {slide_content}")

        return formatted_slides

class RawTextLoader:
    def __init__(self, raw_text: str, metadata={ 'image':'' }):
        self.data = raw_text
        self.metadata = metadata if not None else { 'image': '' }

    def load(self):
        return [RaiLoaderDocument(page_content=self.data, metadata=self.metadata)]

class JSONLLoader:
    def __init__(self, file_path: str, metadata={ 'image':'' }):
        self.file_path = file_path
        self.metadata = metadata if not None else { 'image': '' }

    def load(self):
        docs = []
        try:
            with jsonlines.open(self.file_path) as reader:
                for obj in reader:
                    data_string = DICT.get_any(keys=["content", "text", "page", "page_content"], dic=obj)
                    docs.append(RaiLoaderDocument(page_content=data_string, metadata=self.metadata))
        except Exception as e:
            Log.e(e)
            docs = []
        return docs

    @staticmethod
    def load_file(file:str):
        objs = []
        if os.path.exists(file):
            with open(file, 'r', encoding='utf-8') as f:
                for line in f:
                    try:
                        item = json.loads(line)
                        objs.append(item)
                    except json.JSONDecodeError:
                        continue
            return objs
        else: return objs
    @staticmethod
    def save_file(data: dict, output_file:str):
        with open(output_file, 'a', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False)
            f.write('\n')

class JSONLoader:
    def __init__(self, file_path: str, metadata={ 'image':'' }):
        self.file_path = file_path
        self.metadata = metadata if not None else {'image': ''}

    def load(self):
        with open(self.file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        data_string = DICT.get_any(keys={"content", "text", "page", "page_content"}, dic=data)
        return RaiLoaderDocument(page_content=data_string, metadata=self.metadata)

class VisionDataLoader:

    def __init__(self, file_path: str, metadata:dict={ 'image': '' }, delete_cache=True):
        self.file_path = file_path
        self.metadata = metadata if not None else {'image': ''}
        self.data = VisionExtractor(self.file_path).extract(save_images=True)
        if delete_cache:
            self.delete_cached_folder()
            time.sleep(1)

    def delete_cached_folder(self):
        try:
            out = os.path.splitext(self.file_path)[0]
            Log.w(f"Deleting Cached Output Dir: [ {out}_output ]")
            delete_folder(f"{out}_output")
        except Exception as e:
            Log.e(f"Error removing OS: {e}")

    def load(self, load_text=True, load_images=True):
        Log.w("VisionDataLoader [ load() ] has been called.")
        if not load_text and not load_images:
            return []

        filtered_data = []
        for page in self.data:
            content = ""
            if load_text:
                content = page.get('page_content', '')
            if load_images:
                self.metadata['image'] = page.get('image', '')
            filtered_data.append(RaiLoaderDocument(page_content=content, metadata=self.metadata))
        return filtered_data

    def get_pages(self, load_text=True, load_images=True):
        if self.data is None:
            self.load()
        return self.load(load_text, load_images)

    def get_text(self):
        if self.data is None:
            self.load()
        return [page.get('page_content', '') for page in self.data]

    def get_images(self):
        if self.data is None:
            self.load()
        return [page.get('image', '') for page in self.data]

    def get_images_and_text(self):
        if self.data is None:
            self.load()
        return self.data

    def is_empty(self):
        if self.data is None:
            self.load()
        if not self.data:
            return True
        return not bool(self.data)

    def load_d(self):
        if self.data is None:
            self.load()
        documents = []
        for page in self.data:
            page_content = page.get('page_content', '')
            metadata = {'image': page.get('image', '')}
            documents.append(RaiLoaderDocument(page_content=page_content, metadata=metadata))
        return documents

def delete_folder(folder_path):
    """
    Deletes a folder and all its contents.

    Parameters:
    - folder_path (str): The path to the folder to delete.

    Raises:
    - FileNotFoundError: If the folder does not exist.
    - NotADirectoryError: If the specified path is not a directory.
    - PermissionError: If the operation lacks the necessary permissions.
    - OSError: For other OS-related errors.
    """

    # You can configure logging to output to a file or console as needed

    # Verify that the path exists
    if not os.path.exists(folder_path):
        Log.e(f"Folder not found: {folder_path}")
        raise FileNotFoundError(f"Folder not found: {folder_path}")

    # Verify that the path is a directory
    if not os.path.isdir(folder_path):
        Log.e(f"Specified path is not a directory: {folder_path}")
        raise NotADirectoryError(f"Specified path is not a directory: {folder_path}")

    # Handle read-only files on Windows
    def onerror(func, path, exc_info):
        """
        Error handler for shutil.rmtree.

        If the error is due to an access error (read-only file), it attempts to add write permission and retries.
        If the error is for another reason, it re-raises the error.
        """
        if not os.access(path, os.W_OK):
            # Attempt to add write permission
            os.chmod(path, stat.S_IWUSR)
            func(path)
        else:
            raise

    try:
        shutil.rmtree(folder_path, onerror=onerror)
        Log.i(f"Successfully deleted folder and its contents: {folder_path}")
    except Exception as e:
        Log.e(f"Error deleting folder {folder_path}: {e}")
        raise


known_source_ext = [
            "go",
            "py",
            "java",
            "sh",
            "bat",
            "ps1",
            "cmd",
            "js",
            "ts",
            "css",
            "cpp",
            "hpp",
            "h",
            "c",
            "cs",
            "sql",
            "log",
            "ini",
            "pl",
            "pm",
            "r",
            "dart",
            "dockerfile",
            "env",
            "php",
            "hs",
            "hsc",
            "lua",
            "nginxconf",
            "conf",
            "m",
            "mm",
            "plsql",
            "perl",
            "rb",
            "rs",
            "db2",
            "scala",
            "bash",
            "swift",
            "vue",
            "svelte",
            "msg",
            "ex",
            "exs",
            "erl",
            "tsx",
            "jsx",
            "hs",
            "lhs",
        ]