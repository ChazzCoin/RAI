import os
import csv
from F import LIST
from rai.data.intake.PDF_v1 import FPDF
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from F.LOG import Log
from rai.data.intake.Vision import VisionExtractor

# Configure logging
Log = Log("Files.Read")

""" MASTER FILE OPENER """
def read_file(file_path, enable_vision: bool = False):
    """
    Opens a file and returns its contents as a string, supporting various file formats.
    Supported file types: .txt, .pdf, .docx, .xlsx, .pptx, .csv.
    :param file_path: The path to the file.
    :return: The file contents as a string.
    :raises ValueError: If the file format is unsupported or reading fails.
    """
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_ext = os.path.splitext(file_path)[1].lower()

    if enable_vision:
        vision = __vision_text_extractor(file_path)
        if vision: return vision
        else: Log.w("Vision/OCR Failed, reverting back to basic extraction.")
    try:
        if file_ext == '.txt':
            return _read_txt(file_path)
        elif file_ext == '.pdf':
            return _read_pdf(file_path)
        elif file_ext == '.docx':
            return _read_docx(file_path)
        elif file_ext == '.xlsx':
            return _read_xlsx(file_path)
        elif file_ext == '.pptx':
            return _read_pptx(file_path)
        elif file_ext == '.csv':
            return _read_csv(file_path)
        else:
            Log.e(f"Unsupported file format: {file_ext}")
            return "No Content"

    except Exception as e:
        Log.e(f"Error reading file {file_path}: {e}")
        return "No Content"

# Helper functions for different file formats

def _read_txt(file_path):
    """Reads and returns the contents of a .txt file."""
    Log.i(f"Reading TXT file: {file_path}")
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def __vision_text_extractor(file_path):
    Log.i(f"Attempting Vision OCR Text Extraction on: {file_path}")
    text = []
    try:
        temp = VisionExtractor._vision_extract_text(file_path, save_images=False)
        for item in temp:
            tempText = item["text"]
            text.append(tempText)
        if not text or len(text) == 0:
            return None
        Log.s(f"Vision OCR Text Extraction Success: Page Count: [ {len(text)} ]")
        return '\n'.join(text)
    except Exception as e:
        Log.e(f"Error reading PDF file: {file_path}: {e}")
        return None

def _read_pdf(file_path):
    """Reads and returns the text contents of a .pdf file."""
    Log.i(f"Reading PDF file: {file_path}")
    vision = __vision_text_extractor(file_path)
    if vision:
        return vision
    with open(file_path, 'rb') as file:
        try:
            return LIST.to_str(FPDF.extract_text_from_pdf(file))
        except Exception as e:
            Log.w(f"Error reading PDF Pages. Switching to Fallback... {file_path}: {e}")
    return ""


def _read_docx(file_path):
    """Reads and returns the text contents of a .docx file."""
    Log.i(f"Reading DOCX file: {file_path}")
    vision = __vision_text_extractor(file_path)
    if vision:
        return vision
    doc = DocxDocument(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])


def _read_xlsx(file_path):
    """Reads and returns the contents of an .xlsx file."""
    Log.i(f"Reading XLSX file: {file_path}")
    workbook = load_workbook(file_path, read_only=True)
    contents = []
    for sheet in workbook:
        contents.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            contents.append('\t'.join([str(cell) for cell in row]))
    return '\n'.join(contents)


def _read_pptx(file_path):
    """Reads and returns the text contents of a .pptx file."""
    Log.i(f"Reading PPTX file: {file_path}")
    vision = __vision_text_extractor(file_path)
    if vision:
        return vision
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)


def _read_csv(file_path):
    """Reads and returns the contents of a .csv file."""
    Log.i(f"Reading CSV file: {file_path}")
    with open(file_path, newline='', encoding='utf-8') as csvfile:
        reader = csv.reader(csvfile)
        return '\n'.join([','.join(row) for row in reader])


# Example usage
if __name__ == '__main__':
    file_path = "/path/to/your/file"
    try:
        content = read_file(file_path)
        print(f"Contents of {file_path}:\n{content}")
    except Exception as e:
        print(f"An error occurred: {e}")
